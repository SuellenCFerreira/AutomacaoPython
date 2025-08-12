import tkinter as tk  # Biblioteca para interface gráfica (janelas)
from tkinter import filedialog, messagebox  # Ferramentas para abrir janelas de arquivo e mensagens
import threading  # Para rodar tarefas em paralelo (não travar a janela)
import pandas as pd  # Para manipular planilhas Excel
import openpyxl  # Para manipular arquivos Excel (.xlsx) e desmesclar células
import os  # Para trabalhar com arquivos e pastas
import smtplib  # Para enviar email via protocolo SMTP
from email.message import EmailMessage  # Para criar o email
import mimetypes  # Para identificar tipo do arquivo (anexo)
from pptx import Presentation  # Para criar apresentações PowerPoint
from pptx.util import Inches  # Para posicionar imagens na apresentação
import matplotlib.pyplot as plt  # Para criar gráficos

# ===========================
# Configurações de email Gmail
# ===========================
EMAIL_REMETENTE = 'seu_email@gmail.com'  # Seu email que enviará a mensagem
SENHA_APP = 'sua_senha_app'  # Senha específica do app gerada no Google
SMTP_SERVER = 'smtp.gmail.com'  # Servidor SMTP do Gmail
SMTP_PORT = 587  # Porta usada pelo Gmail para SMTP
EMAIL_DESTINATARIO = 'destino_email@gmail.com'  # Email que vai receber o relatório

# Função para mostrar janela "Aguarde"
def mostrar_aguarde(mensagem="Processando, aguarde..."):
    janela = tk.Toplevel()  # Cria uma nova janela
    janela.title("Aguarde")  # Define o título da janela
    janela.geometry("400x100")  # Define o tamanho da janela (largura x altura)
    janela.resizable(False, False)  # Não permite que a janela seja redimensionada
    tk.Label(janela, text=mensagem, font=("Arial", 12)).pack(expand=True, padx=20, pady=20)  # Coloca texto centralizado
    janela.update()  # Atualiza a janela para aparecer imediatamente
    return janela  # Retorna o objeto da janela para poder fechar depois

# Função para desmesclar células do Excel (openpyxl)
def desmesclar_celulas(caminho_arquivo):
    wb = openpyxl.load_workbook(caminho_arquivo)  # Abre o arquivo Excel para edição
    ws = wb.active  # Seleciona a planilha ativa
    for merged_cell in list(ws.merged_cells.ranges):  # Para cada célula mesclada
        ws.unmerge_cells(str(merged_cell))  # Desfaz a mesclagem
    wb.save(caminho_arquivo)  # Salva as alterações no arquivo

# Função para processar Excel: desmesclar e filtrar colunas
def processar_excel(caminho_entrada, caminho_saida):
    desmesclar_celulas(caminho_entrada)  # Chama a função para desmesclar
    df = pd.read_excel(caminho_entrada)  # Lê o arquivo Excel com pandas
    colunas = ["Data", "Médico", "Motivo", "Status"]  # Define as colunas que queremos manter
    df_filtrado = df[colunas]  # Filtra o DataFrame para só essas colunas
    df_filtrado.to_excel(caminho_saida, index=False)  # Salva o resultado num novo Excel (sem índice)

# Função para gerar gráfico de barras
def gerar_grafico_barras(df, coluna, titulo, caminho_imagem):
    contagem = df[coluna].value_counts()  # Conta quantas vezes cada valor aparece na coluna
    plt.style.use('ggplot')  # Define o estilo do gráfico (visual)
    fig, ax = plt.subplots(figsize=(6,4))  # Cria uma figura e eixos para o gráfico com tamanho
    contagem.plot(kind='bar', ax=ax, color='skyblue')  # Plota o gráfico de barras com cor azul claro
    ax.set_title(titulo)  # Define título do gráfico
    ax.set_ylabel('Quantidade')  # Define o texto do eixo Y
    plt.xticks(rotation=45, ha='right')  # Rotaciona os nomes no eixo X para melhorar a visualização
    plt.tight_layout()  # Ajusta layout para não cortar nada
    plt.savefig(caminho_imagem)  # Salva o gráfico como imagem no caminho definido
    plt.close()  # Fecha o gráfico para liberar memória

# Função para gerar gráfico pizza
def gerar_grafico_pizza(df, coluna, titulo, caminho_imagem):
    contagem = df[coluna].value_counts()  # Conta quantas vezes cada valor aparece na coluna
    plt.style.use('ggplot')  # Define estilo visual
    fig, ax = plt.subplots(figsize=(6,4))  # Cria figura para gráfico
    contagem.plot(kind='pie', ax=ax, autopct='%1.1f%%', startangle=90)  # Gera gráfico pizza com % e começa a desenhar do ângulo 90
    ax.set_ylabel('')  # Remove label do eixo Y (não é necessário)
    ax.set_title(titulo)  # Define título do gráfico
    plt.tight_layout()  # Ajusta layout
    plt.savefig(caminho_imagem)  # Salva imagem do gráfico
    plt.close()  # Fecha o gráfico

# Função para adicionar slide com imagem no PPTX
def adicionar_slide_com_imagem(prs, caminho_imagem, titulo_slide):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Adiciona slide em branco
    if slide.shapes.title:  # Se slide tem título
        slide.shapes.title.text = titulo_slide  # Define texto do título
    left = Inches(1)  # Define margem da esquerda
    top = Inches(1.5)  # Define margem do topo
    slide.shapes.add_picture(caminho_imagem, left, top, width=Inches(6), height=Inches(4))  # Adiciona imagem no slide

# Função para criar apresentação com gráficos
def criar_apresentacao(caminho_excel, caminho_pptx):
    df = pd.read_excel(caminho_excel)  # Lê os dados do Excel processado

    prs = Presentation()  # Cria uma nova apresentação

    # Slide título
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Adiciona slide de título
    slide.shapes.title.text = "Relatório Consultas Médicas"  # Texto título
    slide.placeholders[1].text = "Gerado automaticamente com Python"  # Texto subtítulo

    # Gráfico Consultas por Médico (barra)
    caminho_grafico_medico = "grafico_medico.png"  # Nome do arquivo da imagem do gráfico
    gerar_grafico_barras(df, 'Médico', 'Consultas por Médico', caminho_grafico_medico)  # Cria gráfico barras
    adicionar_slide_com_imagem(prs, caminho_grafico_medico, "Consultas por Médico")  # Adiciona slide com gráfico

    # Gráfico Distribuição por Status (pizza)
    caminho_grafico_status = "grafico_status.png"
    gerar_grafico_pizza(df, 'Status', 'Distribuição por Status da Consulta', caminho_grafico_status)
    adicionar_slide_com_imagem(prs, caminho_grafico_status, "Status das Consultas")

    # Gráfico Motivo das Consultas (barra)
    caminho_grafico_motivo = "grafico_motivo.png"
    gerar_grafico_barras(df, 'Motivo', 'Consultas por Motivo', caminho_grafico_motivo)
    adicionar_slide_com_imagem(prs, caminho_grafico_motivo, "Motivo das Consultas")

    prs.save(caminho_pptx)  # Salva a apresentação no caminho definido

    # Remove os arquivos de imagens temporárias criadas para os gráficos
    for f in [caminho_grafico_medico, caminho_grafico_status, caminho_grafico_motivo]:
        if os.path.exists(f):
            os.remove(f)

# Função para enviar email (em thread para não travar janela)
def enviar_email(arquivo_anexo, assunto, corpo, janela_aguarde=None, destinatario=None):
    try:
        msg = EmailMessage()  # Cria a mensagem de email
        msg['From'] = EMAIL_REMETENTE  # Define quem está enviando
        msg['To'] = destinatario if destinatario else EMAIL_REMETENTE  # Destinatário do email
        msg['Subject'] = assunto  # Assunto do email
        msg.set_content(corpo)  # Corpo do email (texto simples)

        # Descobre o tipo do arquivo para anexar corretamente
        mime_type, _ = mimetypes.guess_type(arquivo_anexo)
        mime_type = mime_type or 'application/octet-stream'
        tipo, sub_tipo = mime_type.split('/')
        with open(arquivo_anexo, 'rb') as f:  # Abre o arquivo para anexar
            msg.add_attachment(f.read(), maintype=tipo, subtype=sub_tipo, filename=os.path.basename(arquivo_anexo))

        # Configura conexão SMTP segura
        with smtplib.SMTP(SMTP_SERVER, SMTP_PORT) as smtp:
            smtp.starttls()  # Inicia criptografia
            smtp.login(EMAIL_REMETENTE, SENHA_APP)  # Login no servidor
            smtp.send_message(msg)  # Envia o email

        # Mostra mensagem de sucesso para o usuário
        messagebox.showinfo("Sucesso", f"E-mail enviado com sucesso para {msg['To']}")

    except Exception as e:
        # Se deu erro, mostra mensagem de erro
        messagebox.showerror("Erro ao enviar e-mail", str(e))

    finally:
        if janela_aguarde:
            janela_aguarde.destroy()  # Fecha janela de aguarde, se estiver aberta

# Função principal que integra tudo
def main():
    root = tk.Tk()  # Cria janela principal
    root.withdraw()  # Esconde essa janela porque não queremos mostrar a raiz

    # Abre janela para escolher arquivo Excel original
    caminho_entrada = filedialog.askopenfilename(
        title="Selecione o arquivo Excel original",
        filetypes=[("Arquivos Excel", "*.xls *.xlsx")])
    if not caminho_entrada:  # Se não escolheu nenhum arquivo, cancela
        messagebox.showinfo("Cancelado", "Nenhum arquivo selecionado.")
        return

    # Abre janela para escolher onde salvar Excel processado
    caminho_saida = filedialog.asksaveasfilename(
        title="Salvar Excel processado como", defaultextension=".xlsx",
        filetypes=[("Excel files", "*.xlsx")])
    if not caminho_saida:  # Se cancelar, sai
        messagebox.showinfo("Cancelado", "Nenhum local para salvar selecionado.")
        return

    # Mostra janela "Aguarde" e processa o Excel (desmescla e filtra)
    aguarde1 = mostrar_aguarde("Processando arquivo Excel, aguarde...")
    aguarde1.update()
    processar_excel(caminho_entrada, caminho_saida)
    aguarde1.destroy()  # Fecha janela aguarde

    messagebox.showinfo("Sucesso", f"Arquivo Excel processado salvo em:\n{caminho_saida}")

    # Abre janela para escolher onde salvar apresentação PPTX
    caminho_pptx = filedialog.asksaveasfilename(
        title="Salvar apresentação PowerPoint como", defaultextension=".pptx",
        filetypes=[("PowerPoint files", "*.pptx")])
    if not caminho_pptx:  # Se cancelar, sai
        messagebox.showinfo("Cancelado", "Nenhum local para salvar a apresentação selecionado.")
        return

    # Mostra janela "Aguarde" e cria a apresentação PowerPoint com gráficos
    aguarde2 = mostrar_aguarde("Criando apresentação, aguarde...")
    aguarde2.update()
    criar_apresentacao(caminho_saida, caminho_pptx)
    aguarde2.destroy()  # Fecha janela aguarde

    messagebox.showinfo("Sucesso", f"Apresentação salva em:\n{caminho_pptx}")

    # Pergunta se quer enviar email para destinatário padrão
    if messagebox.askyesno("Enviar e-mail", "Deseja enviar o relatório por e-mail?"):
        email_destino = EMAIL_DESTINATARIO  # Usa email padrão configurado

        # Mostra janela "Aguarde" para envio
        aguarde_email = mostrar_aguarde("Enviando e-mail, por favor aguarde...")
        aguarde_email.update()

        # Cria thread para enviar email e não travar interface
        def envio_thread():
            enviar_email(
                caminho_pptx,
                "Relatório Consultas Médicas",
                "Segue em anexo o relatório de consultas médicas automatizado.",
                aguarde_email,
                email_destino
            )

        threading.Thread(target=envio_thread).start()

    root.mainloop()  # Mantém o programa aberto até o usuário fechar

# Executa a função principal ao rodar o script
if __name__ == "__main__":
    main()
